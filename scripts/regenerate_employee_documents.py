#!/usr/bin/env python3
"""
Regenerate employee documents with realistic multi-page content.

This script walks data/employees and overwrites each existing file with
format-specific content while preserving file names and relative paths.
"""

from __future__ import annotations

import csv
import random
from datetime import datetime, timedelta, timezone
from email.message import EmailMessage
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


REPO_ROOT = Path(__file__).resolve().parents[1]
EMPLOYEE_DIR = REPO_ROOT / "data" / "employees"
RNG = random.Random(20260511)


TOPICS = [
    "Quarterly delivery milestones",
    "Knowledge transfer and onboarding",
    "Operational risk review",
    "Customer escalation analysis",
    "Automation backlog refinement",
    "Fabric semantic model planning",
    "Data quality remediation",
    "Department KPI performance",
]


def employee_from_path(path: Path) -> str:
    for part in path.parts:
        if part.startswith("EMP"):
            return part
    return "EMP000"


def build_paragraphs(emp_id: str, title: str, section_count: int = 8) -> list[str]:
    paragraphs: list[str] = []
    for idx in range(1, section_count + 1):
        topic = RNG.choice(TOPICS)
        paragraphs.append(
            (
                f"Section {idx}: {topic}. "
                f"{emp_id} contributes to {title.lower()} through recurring collaboration, "
                "clear ownership boundaries, and measurable outcomes. "
                "The team captures decisions, open issues, and handoffs so that knowledge "
                "remains discoverable for future programs."
            )
        )
    return paragraphs


def write_docx(path: Path, emp_id: str, title: str) -> None:
    doc = Document()
    doc.add_heading(f"{title} - {emp_id}", level=1)
    pages = 3
    for page in range(1, pages + 1):
        doc.add_heading(f"Page {page}", level=2)
        for para in build_paragraphs(emp_id, title, section_count=4):
            doc.add_paragraph(para)
        if page < pages:
            doc.add_page_break()
    doc.save(path)


def write_pdf(path: Path, emp_id: str, title: str) -> None:
    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    pages = 3
    for page in range(1, pages + 1):
        y = height - 72
        c.setFont("Helvetica-Bold", 14)
        c.drawString(72, y, f"{title} - {emp_id} (Page {page})")
        y -= 28
        c.setFont("Helvetica", 10)
        for para in build_paragraphs(emp_id, title, section_count=7):
            words = para.split()
            line = []
            for word in words:
                candidate = " ".join(line + [word])
                if c.stringWidth(candidate, "Helvetica", 10) > (width - 144):
                    c.drawString(72, y, " ".join(line))
                    y -= 14
                    line = [word]
                else:
                    line.append(word)
            if line:
                c.drawString(72, y, " ".join(line))
                y -= 18
            if y < 100:
                break
        c.showPage()
    c.save()


def write_pptx(path: Path, emp_id: str, title: str) -> None:
    prs = Presentation()
    for slide_no in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{title} - {emp_id}"
        body = slide.shapes.placeholders[1].text_frame
        body.text = f"Slide {slide_no}: Executive Summary"
        for point in build_paragraphs(emp_id, title, section_count=3):
            p = body.add_paragraph()
            p.text = point
            p.level = 1
    prs.save(path)


def write_xlsx(path: Path, emp_id: str, title: str) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["EmployeeId", "Title", "Week", "Metric", "Value", "Narrative"])
    for week in range(1, 25):
        metric = RNG.choice(["Velocity", "Quality", "Mentoring", "Innovation", "Support"])
        value = RNG.randint(60, 99)
        narrative = build_paragraphs(emp_id, title, section_count=1)[0]
        ws.append([emp_id, title, f"2026-W{week:02d}", metric, value, narrative])

    detail = wb.create_sheet("Actions")
    detail.append(["Date", "Action", "Owner", "Status", "Notes"])
    start = datetime(2026, 1, 1)
    for i in range(1, 41):
        action_date = start + timedelta(days=i * 3)
        detail.append(
            [
                action_date.strftime("%Y-%m-%d"),
                RNG.choice(TOPICS),
                emp_id,
                RNG.choice(["Open", "In Progress", "Closed"]),
                build_paragraphs(emp_id, title, section_count=1)[0],
            ]
        )

    wb.save(path)


def write_csv(path: Path, emp_id: str, title: str) -> None:
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["employeeId", "title", "record", "topic", "status", "summary"])
        for i in range(1, 101):
            writer.writerow(
                [
                    emp_id,
                    title,
                    i,
                    RNG.choice(TOPICS),
                    RNG.choice(["planned", "active", "done"]),
                    build_paragraphs(emp_id, title, section_count=1)[0],
                ]
            )


def write_text_like(path: Path, emp_id: str, title: str, ext: str) -> None:
    pages = []
    for page in range(1, 5):
        header = f"{title} - {emp_id} - Page {page}"
        body = "\n\n".join(build_paragraphs(emp_id, title, section_count=5))
        if ext == ".md":
            pages.append(f"## {header}\n\n{body}")
        else:
            pages.append(f"{header}\n{'=' * len(header)}\n\n{body}")
    path.write_text("\n\n---\n\n".join(pages), encoding="utf-8")


def write_eml(path: Path, emp_id: str, title: str) -> None:
    msg = EmailMessage()
    msg["From"] = f"{emp_id.lower()}@fabrikiq.example"
    msg["To"] = "team@fabrikiq.example"
    msg["Subject"] = f"{title} - Weekly Knowledge Update"
    msg["Date"] = datetime.now(timezone.utc).strftime("%a, %d %b %Y %H:%M:%S +0000")
    body = []
    for idx, para in enumerate(build_paragraphs(emp_id, title, section_count=12), start=1):
        body.append(f"{idx}. {para}")
    msg.set_content("\n\n".join(body))
    path.write_bytes(msg.as_bytes())


def regenerate_file(path: Path) -> str:
    emp_id = employee_from_path(path)
    title = path.stem.replace("_", " ").replace("-", " ").strip().title() or "Knowledge Asset"
    ext = path.suffix.lower()

    if ext == ".docx":
        write_docx(path, emp_id, title)
    elif ext == ".pdf":
        write_pdf(path, emp_id, title)
    elif ext == ".pptx":
        write_pptx(path, emp_id, title)
    elif ext == ".xlsx":
        write_xlsx(path, emp_id, title)
    elif ext == ".csv":
        write_csv(path, emp_id, title)
    elif ext in {".txt", ".md", ".one"}:
        write_text_like(path, emp_id, title, ext)
    elif ext == ".eml":
        write_eml(path, emp_id, title)
    else:
        return "skipped"

    return "updated"


def main() -> int:
    updated = 0
    skipped = 0
    total = 0

    if not EMPLOYEE_DIR.exists():
        raise RuntimeError(f"Employee folder not found: {EMPLOYEE_DIR}")

    for p in EMPLOYEE_DIR.rglob("*"):
        if not p.is_file():
            continue
        total += 1
        result = regenerate_file(p)
        if result == "updated":
            updated += 1
        else:
            skipped += 1

    print("Employee document regeneration complete")
    print(f"  Files scanned: {total}")
    print(f"  Files updated: {updated}")
    print(f"  Files skipped: {skipped}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())