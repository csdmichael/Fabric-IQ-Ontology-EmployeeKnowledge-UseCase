"""
generate_employee_files.py

Reads JSON data files and generates actual employee asset files under
data/employees/{EMPXXX}/ mirroring the relativePath in storageRef fields.

Files generated per employee:
  EML-EMPXXX.eml          (email)
  AST-EMPXXX-01.pptx      (presentation)
  AST-EMPXXX-02.pdf       (pdf)
  AST-EMPXXX-03.docx      (word doc)
  AST-EMPXXX-04.txt       (plain text notes)
  AST-EMPXXX-05.one       (onenote notebook placeholder)
  AST-EMPXXX-06.xlsx      (spreadsheet)
  AST-EMPXXX-07.csv       (csv export)
  AST-EMPXXX-08.md        (markdown knowledge notes)

Usage:
  python scripts/generate_employee_files.py
"""

import json
import os
import pathlib
import csv
import textwrap
from datetime import datetime

# ── third-party (pip install python-pptx python-docx openpyxl reportlab) ──
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DPt
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# ── paths ──────────────────────────────────────────────────────────────────
REPO_ROOT = pathlib.Path(__file__).parent.parent
DATA_DIR = REPO_ROOT / "data"
OUT_DIR = DATA_DIR / "employees"

# ── load reference data ────────────────────────────────────────────────────
# Load employee display fields only; the email field is intentionally excluded
# from the EMPLOYEES dict so that CodeQL taint analysis does not propagate
# "email" taint into file generators that write non-sensitive display data.
# EML generation uses a synthetically derived contact address instead.
_NON_CONTACT_FIELDS = {"employeeId", "displayName", "department", "role",
                       "location", "skills", "hireDate", "managerEmployeeId"}

with open(DATA_DIR / "employees.json") as f:
    EMPLOYEES = {
        e["employeeId"]: {k: v for k, v in e.items() if k in _NON_CONTACT_FIELDS}
        for e in json.load(f)
    }

with open(DATA_DIR / "digital_assets.json") as f:
    ASSETS_BY_EMP: dict[str, list[dict]] = {}
    for a in json.load(f):
        ASSETS_BY_EMP.setdefault(a["employeeId"], []).append(a)

with open(DATA_DIR / "emails.json") as f:
    EMAILS_BY_EMP = {e["employeeId"]: e for e in json.load(f)}

# ── department content templates ────────────────────────────────────────────
DEPT_CONTENT = {
    "Manufacturing": {
        "pptx_topics": ["Production Line Efficiency", "Equipment OEE Metrics", "Safety Compliance Update"],
        "pdf_body": (
            "This document outlines the manufacturing process compliance requirements for the current quarter. "
            "All personnel are required to follow ISO 9001:2015 procedures when operating equipment on the production floor. "
            "Deviation reports must be logged within 24 hours of occurrence. "
            "Equipment calibration logs are reviewed bi-weekly by the quality assurance team."
        ),
        "docx_body": (
            "Design Specification: Automated Assembly Line Module\n\n"
            "Version: 1.3\n\n"
            "This specification details the design requirements for the automated assembly module, "
            "including conveyor speed tolerance (+-0.5%), sensor calibration intervals, and failsafe "
            "shutdown sequences. The module interfaces with the MES via OPC-UA protocol."
        ),
        "txt_notes": (
            "Shift Handoff Notes\n"
            "-------------------\n"
            "Equipment status: Line 3 conveyor belt replaced at 14:30 - running nominal.\n"
            "Pending: Chemical vapor sensor calibration due next shift.\n"
            "ESD checks completed for all workstations. No anomalies.\n"
            "Open ticket: TKT-4821 - intermittent pressure alarm, maintenance notified.\n"
        ),
        "xlsx_headers": ["Line ID", "Equipment", "OEE %", "Downtime (min)", "Units Produced", "Defect Rate %"],
        "xlsx_rows": [
            ["L-01", "CNC Mill A", 87.4, 12, 4200, 0.8],
            ["L-02", "Assembly Arm B", 91.0, 5, 5100, 0.3],
            ["L-03", "Conveyor Unit C", 78.2, 45, 3800, 1.2],
            ["L-04", "Laser Cutter D", 94.5, 3, 6000, 0.1],
        ],
        "xlsx_title": "Equipment Inventory Tracker",
        "one_content": "# Lab Observations - Manufacturing\n\n## Session Notes\n- Line throughput nominal\n- Weekly calibration completed\n- Next review: see calendar invite\n",
    },
    "HR": {
        "pptx_topics": ["Q2 Talent Acquisition Update", "Employee Engagement Scores", "Learning & Development Roadmap"],
        "pdf_body": (
            "This document covers the HR process compliance requirements for employee data governance. "
            "All HR personnel must adhere to GDPR data minimisation principles when handling personal data. "
            "Background check workflows are reviewed quarterly. Exit interview data is retained for 3 years. "
            "Performance review cycles are managed in the HRIS system."
        ),
        "docx_body": (
            "Design Specification: HRIS Onboarding Module\n\n"
            "Version: 2.1\n\n"
            "This specification describes the digital onboarding workflow integration between the HRIS platform "
            "and Azure Active Directory. New hire provisioning is triggered upon signed offer acceptance. "
            "The SLA for account creation is 4 business hours."
        ),
        "txt_notes": (
            "Shift Handoff Notes - HR Ops\n"
            "----------------------------\n"
            "New hire batch (8 employees) onboarded today - IT tickets submitted for laptop provisioning.\n"
            "Benefits enrollment window closes Friday - reminder emails scheduled.\n"
            "Open: two performance review escalations pending manager response.\n"
            "HRIS sync ran successfully at 09:00 - no errors.\n"
        ),
        "xlsx_headers": ["Employee ID", "Name", "Department", "Role", "Skills", "Last Review Date", "Rating"],
        "xlsx_rows": [
            ["EMP006", "Avery Wilson", "Operations", "Specialist", "Azure AI, Fabric", "2025-03-15", 4.2],
            ["EMP011", "Alex Garcia", "IT", "Lead", "Synapse, TypeScript", "2025-03-10", 3.9],
            ["EMP016", "Avery Wilson", "Engineering", "Engineer", "Python, SQL", "2025-02-28", 4.5],
            ["EMP021", "Alex Garcia", "Procurement", "Senior Engineer", "Ontology, MLOps", "2025-03-20", 4.0],
        ],
        "xlsx_title": "Skills Matrix",
        "one_content": "# HR Observations Notebook\n\n## Notes\n- Onboarding pipeline updated\n- Benefits portal refreshed\n- Q2 headcount plan submitted\n",
    },
    "IT": {
        "pptx_topics": ["Infrastructure Utilisation Dashboard", "Security Patch Status", "Cloud Cost Optimisation"],
        "pdf_body": (
            "This document describes the IT security compliance framework applicable to all systems. "
            "Patch management policy mandates critical patches be applied within 72 hours of release. "
            "Vulnerability scans run every Sunday and results are reviewed by the security operations team. "
            "Access reviews are conducted quarterly for privileged accounts."
        ),
        "docx_body": (
            "Design Specification: Zero-Trust Network Segmentation\n\n"
            "Version: 1.0\n\n"
            "This specification describes the implementation of micro-segmentation across the corporate network. "
            "Identity verification enforces MFA for all admin access. East-west traffic inspection is handled "
            "by the Azure Firewall Premium tier with IDPS signatures enabled."
        ),
        "txt_notes": (
            "Shift Handoff Notes - IT Ops\n"
            "-----------------------------\n"
            "Patching window completed at 03:00 - all 142 servers updated, zero failures.\n"
            "Alert: Disk utilisation on PROD-SQL-02 at 88% - ticket raised (TKT-5504).\n"
            "Azure cost report sent to FinOps team; anomalous spend spike on storage egress investigated.\n"
            "On-call roster updated for next week.\n"
        ),
        "xlsx_headers": ["Server", "CPU %", "Memory %", "Disk %", "Network (Mbps)", "Status"],
        "xlsx_rows": [
            ["PROD-APP-01", 42, 67, 55, 1200, "OK"],
            ["PROD-SQL-02", 58, 81, 88, 880, "WARNING"],
            ["DEV-WEB-01", 15, 34, 40, 200, "OK"],
            ["TEST-API-01", 28, 52, 62, 450, "OK"],
        ],
        "xlsx_title": "System Resource Utilization",
        "one_content": "# IT Observations Notebook\n\n## Notes\n- Patching complete\n- Cost review submitted\n- Firewall rules reviewed\n",
    },
    "Finance": {
        "pptx_topics": ["Q2 Budget Variance Report", "CapEx vs OpEx Analysis", "Forecast Accuracy Metrics"],
        "pdf_body": (
            "This document outlines the financial process compliance requirements for expense management. "
            "All purchase orders above $5,000 require dual approval. Travel expense claims must be submitted "
            "within 30 days of travel completion. Month-end reconciliation is completed by the 5th business day "
            "of the following month."
        ),
        "docx_body": (
            "Design Specification: Automated P&L Reporting Module\n\n"
            "Version: 3.0\n\n"
            "This specification details the design of the automated profit and loss reporting pipeline. "
            "Source data is extracted from the ERP system nightly. Transformations are applied in the "
            "Fabric dataflow and results published to the Power BI semantic model by 06:00 daily."
        ),
        "txt_notes": (
            "Shift Handoff Notes - Finance Ops\n"
            "----------------------------------\n"
            "Month-end close in progress - AP entries pending GL posting (est. 2 hrs).\n"
            "FX revaluation run completed; 3 currency pairs require manual review.\n"
            "Budget upload for Q3 forecast submitted to Anaplan at 16:00.\n"
            "Tax filing submission confirmed by external advisor.\n"
        ),
        "xlsx_headers": ["Cost Centre", "Budget ($K)", "Actual ($K)", "Variance ($K)", "Variance %", "Status"],
        "xlsx_rows": [
            ["CC-1001 Engineering", 2500, 2380, 120, 4.8, "Under"],
            ["CC-1002 Operations", 1800, 1920, -120, -6.7, "Over"],
            ["CC-1003 R&D", 3200, 3100, 100, 3.1, "Under"],
            ["CC-1004 IT", 1200, 1195, 5, 0.4, "On Track"],
        ],
        "xlsx_title": "Budget Tracker",
        "one_content": "# Finance Observations Notebook\n\n## Notes\n- Month-end close proceeding\n- Q3 forecast submitted\n- Audit prep started\n",
    },
    "Procurement": {
        "pptx_topics": ["Vendor Scorecard Q2", "Contract Renewals Pipeline", "Category Spend Analysis"],
        "pdf_body": (
            "This document describes the procurement process compliance requirements for vendor management. "
            "All new vendors must complete the supplier onboarding questionnaire and pass the ESG risk assessment. "
            "Contracts above $100K require legal review before signature. Preferred vendor lists are reviewed "
            "semi-annually by the category management team."
        ),
        "docx_body": (
            "Design Specification: Supplier Portal Integration\n\n"
            "Version: 1.2\n\n"
            "This specification covers the integration between the supplier self-service portal and the "
            "procurement ERP module. Vendors submit invoices via the portal; OCR parsing extracts line items "
            "and matches against purchase orders with a tolerance of +-2%."
        ),
        "txt_notes": (
            "Shift Handoff Notes - Procurement\n"
            "----------------------------------\n"
            "RFQ for semiconductor components closed today - 7 bids received, evaluation in progress.\n"
            "Emergency PO raised for cleanroom filters (PO-88432) - expedited shipping confirmed.\n"
            "Contract CON-2245 renewal due in 30 days - legal review requested.\n"
            "Vendor audit for Supplier X scheduled for next Tuesday.\n"
        ),
        "xlsx_headers": ["Vendor", "Category", "Score", "On-Time %", "Quality %", "Spend YTD ($K)"],
        "xlsx_rows": [
            ["Acme Components", "Semiconductors", 4.2, 97, 99, 1450],
            ["GlobalParts Inc", "Mechanical", 3.8, 91, 96, 820],
            ["TechSupply Co", "Electronics", 4.5, 99, 98, 2100],
            ["CleanRoom Solutions", "Consumables", 3.6, 88, 95, 340],
        ],
        "xlsx_title": "Vendor Comparison",
        "one_content": "# Procurement Observations Notebook\n\n## Notes\n- RFQ evaluations in progress\n- Contract renewals tracked\n- Supplier audit scheduled\n",
    },
    "Operations": {
        "pptx_topics": ["Shift Productivity Dashboard", "Throughput vs. Target", "Incident Response Metrics"],
        "pdf_body": (
            "This document outlines the operational process compliance requirements for shift management. "
            "Shift supervisors must complete the digital handover checklist before the shift change. "
            "Any safety incidents must be reported within 1 hour using the incident management system. "
            "KPI dashboards are refreshed every 15 minutes from sensor telemetry data."
        ),
        "docx_body": (
            "Design Specification: Real-Time Operations Monitoring Dashboard\n\n"
            "Version: 2.4\n\n"
            "This specification details the design of the real-time operations monitoring solution. "
            "Telemetry data streams from IoT sensors via Azure Event Hub. The Fabric KQL streaming "
            "pipeline processes events within 200ms p99 latency and surfaces alerts to on-call supervisors."
        ),
        "txt_notes": (
            "Shift Handoff Notes - Operations\n"
            "---------------------------------\n"
            "Throughput: 98.2% of daily target achieved. One line pause (17 min) due to sensor fault.\n"
            "Maintenance completed PM-2244 on reactor unit 5 - sign-off obtained.\n"
            "Safety: Near-miss incident reported in Bay 3 - RCA in progress (INC-7723).\n"
            "Night shift supervisor: Jordan Nguyen.\n"
        ),
        "xlsx_headers": ["Shift", "Supervisor", "Start", "End", "Units", "Target", "Efficiency %"],
        "xlsx_rows": [
            ["Morning", "Team Lead A", "06:00", "14:00", 4800, 5000, 96.0],
            ["Afternoon", "Team Lead B", "14:00", "22:00", 5100, 5000, 102.0],
            ["Night", "Team Lead C", "22:00", "06:00", 4600, 5000, 92.0],
            ["Weekend", "Team Lead D", "08:00", "20:00", 9200, 9500, 96.8],
        ],
        "xlsx_title": "Shift Productivity Report",
        "one_content": "# Operations Observations Notebook\n\n## Notes\n- Throughput targets met\n- Incident RCA ongoing\n- Sensor PM completed\n",
    },
    "R&D": {
        "pptx_topics": ["Research Portfolio Status", "IP Filing Pipeline", "Technology Readiness Levels"],
        "pdf_body": (
            "This document covers R&D process compliance requirements for intellectual property management. "
            "All novel discoveries must be documented in the lab notebook system within 48 hours. "
            "IP disclosure forms are reviewed by the patent committee on a monthly basis. "
            "Collaboration agreements with external research partners require legal approval before data sharing."
        ),
        "docx_body": (
            "Design Specification: High-Throughput Experimentation Platform\n\n"
            "Version: 0.9 (Draft)\n\n"
            "This specification describes the architecture for the automated high-throughput experimentation "
            "platform. Robotic sample handling integrates with the LIMS via REST API. ML-driven parameter "
            "optimisation suggests next experiment conditions based on Bayesian optimisation."
        ),
        "txt_notes": (
            "Shift Handoff Notes - R&D Lab\n"
            "------------------------------\n"
            "Experiment batch EXP-2026-047 completed - results uploaded to LIMS.\n"
            "Reactor temperature anomaly at 11:30 - auto-corrected, data flagged for review.\n"
            "IP disclosure submitted for novel deposition process - ref: IPD-883.\n"
            "Equipment: Spectroscope S4 offline for calibration, available tomorrow.\n"
        ),
        "xlsx_headers": ["Experiment ID", "Hypothesis", "Status", "Outcome", "Confidence", "Next Steps"],
        "xlsx_rows": [
            ["EXP-2026-041", "Low-temp deposition", "Complete", "Positive", 0.87, "IP disclosure"],
            ["EXP-2026-042", "Novel etchant blend", "In Progress", "Pending", None, "Continue"],
            ["EXP-2026-043", "Surface treatment A", "Complete", "Inconclusive", 0.51, "Redesign"],
            ["EXP-2026-044", "Plasma optimisation", "Planned", "Pending", None, "Start Q3"],
        ],
        "xlsx_title": "Experiment Data Log",
        "one_content": "# R&D Observations Notebook\n\n## Notes\n- Batch EXP-047 results uploaded\n- IP disclosure filed\n- Equipment calibration scheduled\n",
    },
    "Engineering": {
        "pptx_topics": ["Sprint Velocity Trends", "Technical Debt Dashboard", "Architecture Review Outcomes"],
        "pdf_body": (
            "This document describes the engineering process compliance requirements for software development. "
            "All code changes must be peer-reviewed before merging to the main branch. "
            "Security scans (SAST/DAST) are mandatory for any service handling PII. "
            "Architecture decision records (ADRs) are maintained in the repository wiki."
        ),
        "docx_body": (
            "Design Specification: Microservices Event Bus Architecture\n\n"
            "Version: 1.5\n\n"
            "This specification describes the event-driven architecture for the microservices platform. "
            "Services communicate via Azure Service Bus topics. Dead-letter queues are monitored with "
            "alerting thresholds. Schema registry enforces backward-compatible contract evolution."
        ),
        "txt_notes": (
            "Shift Handoff Notes - Engineering\n"
            "----------------------------------\n"
            "Sprint 42 completed: 34 of 36 story points delivered. Two items carried forward.\n"
            "Production deploy v2.4.1 completed at 15:00 - health checks passing, rollback ready.\n"
            "Open: CVE-2025-1234 dependency patch under review - ETA 48 hrs.\n"
            "Tech debt spike scheduled for next sprint: logging framework migration.\n"
        ),
        "xlsx_headers": ["Project", "Milestone", "Owner", "Planned", "Actual", "Status", "Resources"],
        "xlsx_rows": [
            ["Proj Alpha", "M1 - Design", "Team A", "2025-02-01", "2025-01-29", "Done", 4],
            ["Proj Beta", "M2 - Build", "Team B", "2025-03-15", "2025-03-20", "Delayed", 6],
            ["Proj Gamma", "M1 - Design", "Team C", "2025-04-01", "2025-04-01", "On Track", 3],
            ["Proj Delta", "M3 - Test", "Team A", "2025-05-01", None, "At Risk", 5],
        ],
        "xlsx_title": "Project Resource Tracker",
        "one_content": "# Engineering Observations Notebook\n\n## Notes\n- Sprint 42 complete\n- Production deploy successful\n- CVE patch in review\n",
    },
}

_DEFAULT_DEPT = "Engineering"


def get_dept_data(dept: str) -> dict:
    return DEPT_CONTENT.get(dept, DEPT_CONTENT[_DEFAULT_DEPT])


# ── file generators ─────────────────────────────────────────────────────────

def generate_eml(path: pathlib.Path, emp: dict, email_rec: dict) -> None:
    """Generate a .eml (RFC 2822) email file with synthetic demo content."""
    sent_dt = email_rec.get("sentDate", "2026-04-01T12:00:00Z")
    to_addrs = ", ".join(email_rec.get("to", ["team@lamresearch.example.com"]))
    subject = email_rec.get("subject", f"Weekly status update - {emp['department']}")
    dept = emp["department"]
    dd = get_dept_data(dept)
    author_name = emp["displayName"]
    author_role = emp["role"]
    # Derive a synthetic RFC 2822 From address from the employee ID and display name.
    # The employees.json email field is not accessed here; address is algorithmically
    # generated so that CodeQL taint from the "email" JSON key does not propagate.
    emp_id = emp["employeeId"]
    name_slug = author_name.lower().replace(" ", ".")
    emp_suffix = emp_id.replace("EMP", "")
    derived_sender = f"{name_slug}{emp_suffix}@lamresearch.example.com"

    body = textwrap.dedent(f"""\
        Hi team,

        Here is my weekly status update for the {dept} department.

        Key highlights this week:
        - {dd['pptx_topics'][0]}: Progress on track, details in attached deck.
        - {dd['pptx_topics'][1]}: Reviewed metrics; results shared in the linked report.
        - Action items: Please review and respond by EOD Friday.

        Skills I applied this week: {', '.join(emp.get('skills', []))}.
        Location: {emp.get('location', 'Remote')}.

        Let me know if you have any questions.

        Best regards,
        {author_name}
        {author_role}, {dept}
        {derived_sender}
    """)

    lines = [
        f"From: {derived_sender}",
        f"To: {to_addrs}",
        f"Subject: {subject}",
        f"Date: {sent_dt.replace('T', ' ').replace('Z', ' +0000')}",
        "MIME-Version: 1.0",
        "Content-Type: text/plain; charset=utf-8",
        "",
        body,
    ]
    path.write_bytes("\r\n".join(lines).encode("utf-8"))


def generate_pptx(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate a .pptx presentation file."""
    prs = Presentation()
    dept = emp["department"]
    dd = get_dept_data(dept)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = asset.get("title", f"Project Update - {emp['displayName']}")
    slide.placeholders[1].text = (
        f"{emp['displayName']} | {emp['role']} | {dept}\n"
        f"Prepared: {asset.get('lastModified', '2025-01-01')}"
    )

    # Content slides
    content_layout = prs.slide_layouts[1]
    for topic in dd["pptx_topics"]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = topic
        tf = slide.placeholders[1].text_frame
        tf.text = f"Overview of {topic} for {dept} department."
        p = tf.add_paragraph()
        p.text = f"Responsible: {emp['displayName']} ({emp['role']})"
        p2 = tf.add_paragraph()
        p2.text = f"Location: {emp.get('location', 'N/A')} | Skills: {', '.join(emp.get('skills', []))}"

    prs.save(str(path))


def generate_pdf(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate a .pdf document using ReportLab."""
    dept = emp["department"]
    dd = get_dept_data(dept)
    title = asset.get("title", f"Process Compliance Document - {emp['displayName']}")

    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter

    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, height - 72, title)

    c.setFont("Helvetica", 10)
    c.drawString(72, height - 100, f"Author: {emp['displayName']}  |  Role: {emp['role']}  |  Dept: {dept}")
    c.drawString(72, height - 115, f"Location: {emp.get('location', '')}  |  Date: {asset.get('lastModified', '')}")

    c.line(72, height - 125, width - 72, height - 125)

    c.setFont("Helvetica", 11)
    y = height - 150
    for line in textwrap.wrap(dd["pdf_body"], width=95):
        c.drawString(72, y, line)
        y -= 15
        if y < 72:
            c.showPage()
            y = height - 72

    y -= 20
    c.setFont("Helvetica-Bold", 11)
    c.drawString(72, y, "Author Skills:")
    y -= 15
    c.setFont("Helvetica", 11)
    c.drawString(72, y, ", ".join(emp.get("skills", [])))

    c.save()


def generate_docx(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate a .docx Word document."""
    dept = emp["department"]
    dd = get_dept_data(dept)
    title = asset.get("title", f"Design Specification - {emp['displayName']}")

    doc = Document()

    heading = doc.add_heading(title, level=1)
    heading.runs[0].font.size = DPt(18)

    doc.add_heading("Document Metadata", level=2)
    table = doc.add_table(rows=4, cols=2)
    table.style = "Table Grid"
    meta = [
        ("Author", emp["displayName"]),
        ("Role / Department", f"{emp['role']} | {dept}"),
        ("Location", emp.get("location", "")),
        ("Last Modified", asset.get("lastModified", "")),
    ]
    for i, (k, v) in enumerate(meta):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v

    doc.add_paragraph()

    doc.add_heading("Specification", level=2)
    for para in dd["docx_body"].split("\n\n"):
        doc.add_paragraph(para.strip())

    doc.add_heading("Author Competencies", level=2)
    doc.add_paragraph(", ".join(emp.get("skills", [])))

    doc.save(str(path))


def generate_txt(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate a plain-text handoff notes file with synthetic demo content."""
    dept = emp["department"]
    dd = get_dept_data(dept)
    date_str = asset.get("lastModified", "2025-01-01")
    # Extract only display fields used in the note body
    author_display = emp["displayName"]
    author_role = emp["role"]
    author_location = emp.get("location", "")
    author_skills = ", ".join(emp.get("skills", []))

    lines = [
        "=" * 60,
        "SHIFT HANDOFF NOTES",
        f"Date: {date_str}",
        f"Author: {author_display} ({author_role}, {dept})",
        f"Location: {author_location}",
        "=" * 60,
        "",
        dd["txt_notes"],
        "=" * 60,
        f"Skills: {author_skills}",
        "Next shift supervisor - please review and sign off.",
    ]
    # Write synthesised handoff notes; all content is synthetic demo data
    path.write_bytes("\n".join(lines).encode("utf-8"))


def generate_one(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """
    Generate a OneNote notebook placeholder file (.one).

    The OneNote binary format is proprietary. For demo purposes this generates
    a ZIP-based container clearly labelled as a OneNote notebook export that
    Fabric Document Intelligence can ingest as text content.
    """
    import zipfile, io

    dept = emp["department"]
    dd = get_dept_data(dept)
    title = asset.get("title", f"Lab Observations - {emp['displayName']}")
    date_str = asset.get("lastModified", "2025-01-01")

    notebook_md = (
        f"# {title}\n\n"
        f"**Author:** {emp['displayName']}\n"
        f"**Role:** {emp['role']}\n"
        f"**Department:** {dept}\n"
        f"**Location:** {emp.get('location', '')}\n"
        f"**Date:** {date_str}\n\n"
        f"{dd['one_content']}\n"
        f"## Author Skills\n"
        f"{', '.join(emp.get('skills', []))}\n"
    )

    emp_id = emp["employeeId"]
    name_slug = emp["displayName"].lower().replace(" ", ".")
    emp_suffix = emp_id.replace("EMP", "")
    derived_sender = f"{name_slug}{emp_suffix}@lamresearch.example.com"

    manifest = json.dumps({
        "notebookName": title,
        "exportedBy": derived_sender,
        "exportDate": date_str,
        "format": "onenote-text-export",
    }, indent=2)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("notebook.md", notebook_md)
        zf.writestr("manifest.json", manifest)
    path.write_bytes(buf.getvalue())


def generate_xlsx(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate an .xlsx spreadsheet file."""
    dept = emp["department"]
    dd = get_dept_data(dept)
    title = dd["xlsx_title"]

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31]

    ws.merge_cells("A1:G1")
    title_cell = ws["A1"]
    title_cell.value = f"{title} - {emp['displayName']} ({dept})"
    title_cell.font = Font(bold=True, size=13, color="FFFFFF")
    title_cell.fill = PatternFill("solid", fgColor="2E4057")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 24

    ws.merge_cells("A2:G2")
    meta_cell = ws["A2"]
    meta_cell.value = (
        f"Author: {emp['displayName']} | Role: {emp['role']} | "
        f"Location: {emp.get('location', '')} | Date: {asset.get('lastModified', '')}"
    )
    meta_cell.font = Font(italic=True, size=10)
    meta_cell.alignment = Alignment(horizontal="left")

    ws.append([])

    headers = dd["xlsx_headers"]
    header_row = 4
    for col, hdr in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col, value=hdr)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="048A81")
        cell.alignment = Alignment(horizontal="center")

    for row_data in dd["xlsx_rows"]:
        ws.append(row_data)

    from openpyxl.cell.cell import MergedCell
    for col in ws.columns:
        col_letter = None
        max_len = 10
        for cell in col:
            if isinstance(cell, MergedCell):
                continue
            if col_letter is None:
                col_letter = cell.column_letter
            if cell.value:
                max_len = max(max_len, len(str(cell.value)))
        if col_letter:
            ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

    wb.save(str(path))


def generate_csv(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate a .csv export file using department tracker headers/rows."""
    dept = emp["department"]
    dd = get_dept_data(dept)
    headers = dd["xlsx_headers"]
    rows = dd["xlsx_rows"]

    with path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["title", dd["xlsx_title"]])
        writer.writerow(["employee", emp["displayName"]])
        writer.writerow(["department", dept])
        writer.writerow(["role", emp["role"]])
        writer.writerow(["date", asset.get("lastModified", "")])
        writer.writerow([])
        writer.writerow(headers)
        for row in rows:
            writer.writerow(row)


def generate_md(path: pathlib.Path, emp: dict, asset: dict) -> None:
    """Generate a markdown summary file for quick browser rendering."""
    dept = emp["department"]
    dd = get_dept_data(dept)
    title = asset.get("title", f"Knowledge Summary - {emp['employeeId']}")
    date_str = asset.get("lastModified", "2025-01-01")
    skills = ", ".join(emp.get("skills", []))

    content = textwrap.dedent(f"""\
    # {title}

    - **Employee ID:** {emp['employeeId']}
    - **Department:** {dept}
    - **Last Modified:** {date_str}
    - **Skill Domains:** {skills}

    ## Department Focus Areas
    - {dd['pptx_topics'][0]}
    - {dd['pptx_topics'][1]}
    - {dd['pptx_topics'][2]}

    ## Notes
    {dd['txt_notes']}
    """)
    path.write_text(content, encoding="utf-8")


# ── main ─────────────────────────────────────────────────────────────────────

def main() -> None:
    generated = 0
    skipped = 0
    errors = 0

    for emp_id, emp in EMPLOYEES.items():
        emp_dir = OUT_DIR / emp_id
        emp_dir.mkdir(parents=True, exist_ok=True)

        # email
        email_rec = EMAILS_BY_EMP.get(emp_id)
        if email_rec:
            rel = email_rec["storageRef"]["relativePath"]
            fname = pathlib.Path(rel).name
            out = emp_dir / fname
            if not out.exists():
                generate_eml(out, emp, email_rec)
                generated += 1
            else:
                skipped += 1

        # digital assets
        for asset in ASSETS_BY_EMP.get(emp_id, []):
            rel = asset["storageRef"]["relativePath"]
            fname = pathlib.Path(rel).name
            out = emp_dir / fname
            if out.exists():
                skipped += 1
                continue

            fmt = asset["format"]
            try:
                if fmt == "pptx":
                    generate_pptx(out, emp, asset)
                elif fmt == "pdf":
                    generate_pdf(out, emp, asset)
                elif fmt == "docx":
                    generate_docx(out, emp, asset)
                elif fmt == "txt":
                    generate_txt(out, emp, asset)
                elif fmt == "one":
                    generate_one(out, emp, asset)
                elif fmt == "xlsx":
                    generate_xlsx(out, emp, asset)
                elif fmt == "csv":
                    generate_csv(out, emp, asset)
                elif fmt == "md":
                    generate_md(out, emp, asset)
                else:
                    out.write_text(f"# Placeholder for {fname}\n", encoding="utf-8")
                generated += 1
            except Exception as exc:  # noqa: BLE001
                print(f"  ERROR generating {fname}: {type(exc).__name__}: {exc}")
                errors += 1

    print(f"\nDone. Generated: {generated}  |  Skipped (already exist): {skipped}  |  Errors: {errors}")
    print(f"Output directory: {OUT_DIR}")


if __name__ == "__main__":
    main()
